import time
import asyncio
import logging
import tempfile
import zipfile
from pathlib import Path
from typing import Dict, List, Optional, Any, Union, Tuple, Iterator
import xml.etree.ElementTree as ET
from concurrent.futures import ThreadPoolExecutor
import base64
import io
from dataclasses import dataclass
from enum import Enum
import re
import subprocess
from utils.auto_cleanup import auto_cleanup

# Modern document processing with Docling
try:
    from docling.document_converter import DocumentConverter
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.document import DoclingDocument
    DOCLING_AVAILABLE = True
except ImportError:
    DOCLING_AVAILABLE = False
    logging.warning("Docling not available. Install with: pip install docling")

# Traditional office libraries
try:
    from docx import Document as DocxDocument
    from docx.shared import Inches, RGBColor
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
    from docx.oxml.ns import qn
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False
    logging.warning("python-docx not available")

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches as PptxInches
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False
    logging.warning("python-pptx not available")

try:
    import openpyxl
    from openpyxl.utils import get_column_letter, column_index_from_string
    from openpyxl.workbook import Workbook
    from openpyxl.worksheet.worksheet import Worksheet
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
    logging.warning("openpyxl not available")

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    from PIL import Image
    import cv2
    import numpy as np
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False
    logging.warning("ReportLab not available. Install with: pip install reportlab")

try:
    import matplotlib.pyplot as plt
    from matplotlib.backends.backend_pdf import PdfPages
    MATPLOTLIB_AVAILABLE = True
except ImportError:
    MATPLOTLIB_AVAILABLE = False
    logging.warning("Matplotlib not available. Install with: pip install matplotlib")


# Base classes
from .base import (
    DocumentProcessor, ProcessingConfig, ProcessingResult, 
    DocumentMetadata, DocumentFormat, ProcessingMode
)


class OfficeContentType(Enum):
    """Types of content that can be extracted from office documents"""
    TEXT = "text"
    TABLE = "table"
    IMAGE = "image"
    CHART = "chart"
    FORMULA = "formula"
    HYPERLINK = "hyperlink"
    HEADER_FOOTER = "header_footer"
    COMMENT = "comment"
    SHAPE = "shape"


@dataclass
class ExtractedContent:
    """Container for extracted office content"""
    content_type: OfficeContentType
    content: str
    metadata: Dict[str, Any]
    position: Optional[Dict[str, Any]] = None
    formatting: Optional[Dict[str, Any]] = None
    confidence: float = 1.0

class DocxToPdfProcessor:
    """Convert DOCX to PDF then process with PDF processors"""
    
    def __init__(self, config):
        self.config = config
    
    def process(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """Process DOCX by converting to PDF first"""
        file_path = Path(file_path)
        
        if not file_path.suffix.lower() == '.docx':
            return {
                "content": "",
                "metadata": {},
                "success": False,
                "error": "Not a DOCX file"
            }
        
        try:
            # Convert DOCX to PDF
            pdf_path = self._convert_to_pdf(file_path)
            
            if not pdf_path or not pdf_path.exists():
                return self._fallback_to_python_docx(file_path)
            
            # Process PDF with our existing PDF processors
            from .pdf_processor import PDFProcessor
            pdf_processor = PDFProcessor(self.config)
            result = pdf_processor.process(pdf_path)
            
            # Clean up temp PDF
            try:
                pdf_path.unlink()
            except:
                pass
            
            if result.success:
                return {
                    "content": result.content,
                    "metadata": {
                        "title": result.metadata.title or file_path.stem,
                        "format": "docx",
                        "processing_engine": "docx-to-pdf-conversion",
                        "original_format": "docx",
                        "converted_via": "pdf"
                    },
                    "images": result.images or [],
                    "tables": result.tables or [],
                    "formulas": result.formulas or [],
                    "success": True,
                    "processor": "docx-to-pdf"
                }
            else:
                # Fallback to direct DOCX processing
                return self._fallback_to_python_docx(file_path)
                
        except Exception as e:
            logging.warning(f"DOCX to PDF conversion failed: {e}")
            return self._fallback_to_python_docx(file_path)
    
    def _convert_to_pdf(self, docx_path: Path) -> Path:
        """Convert DOCX to PDF using available tools"""
        temp_dir = Path(tempfile.gettempdir())
        pdf_path = temp_dir / f"{docx_path.stem}_temp.pdf"

        auto_cleanup(pdf_path, timeout=1800)
        
        # Try multiple conversion methods
        conversion_methods = [
            self._convert_with_libreoffice,
            self._convert_with_unoconv,
            self._convert_with_pandoc,
            self._convert_with_python_docx2pdf
        ]
        
        for method in conversion_methods:
            try:
                if method(docx_path, pdf_path):
                    logging.info(f"Successfully converted {docx_path} to PDF using {method.__name__}")
                    return pdf_path
            except Exception as e:
                logging.debug(f"Conversion method {method.__name__} failed: {e}")
                continue
        
        logging.warning("All PDF conversion methods failed")
        return None
    
    def _convert_with_libreoffice(self, docx_path: Path, pdf_path: Path) -> bool:
        """Convert using LibreOffice headless"""
        try:
            cmd = [
                'libreoffice', '--headless', '--convert-to', 'pdf',
                '--outdir', str(pdf_path.parent),
                str(docx_path)
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
            
            # LibreOffice creates PDF with same name as input
            expected_pdf = pdf_path.parent / f"{docx_path.stem}.pdf"
            if expected_pdf.exists():
                expected_pdf.rename(pdf_path)
                return True
            
            return result.returncode == 0 and pdf_path.exists()
        except:
            return False
    
    def _convert_with_unoconv(self, docx_path: Path, pdf_path: Path) -> bool:
        """Convert using unoconv"""
        try:
            cmd = ['unoconv', '-f', 'pdf', '-o', str(pdf_path), str(docx_path)]
            result = subprocess.run(cmd, capture_output=True, timeout=30)
            return result.returncode == 0 and pdf_path.exists()
        except:
            return False
    
    def _convert_with_pandoc(self, docx_path: Path, pdf_path: Path) -> bool:
        """Convert using Pandoc"""
        try:
            cmd = ['pandoc', str(docx_path), '-o', str(pdf_path)]
            result = subprocess.run(cmd, capture_output=True, timeout=30)
            return result.returncode == 0 and pdf_path.exists()
        except:
            return False
    
    def _convert_with_python_docx2pdf(self, docx_path: Path, pdf_path: Path) -> bool:
        """Convert using python docx2pdf library"""
        try:
            from docx2pdf import convert
            convert(str(docx_path), str(pdf_path))
            return pdf_path.exists()
        except ImportError:
            return False
        except:
            return False
    
    def _fallback_to_python_docx(self, file_path: Path) -> Dict[str, Any]:
        """Fallback to direct python-docx processing"""
        try:
            # FIX: Sửa import path
            processor = EnhancedDocxProcessor(self.config)  # Thay vì from .office_processor import
            return processor.process(file_path)
        except Exception as e:
            return {
                "content": "",
                "metadata": {},
                "success": False,
                "error": f"All processing methods failed: {e}",
                "processor": "docx-fallback"
            }


class PptxToPdfProcessor:
    """Convert PPTX to PDF then process with PDF processors"""
    
    def __init__(self, config):
        self.config = config
    
    def process(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """Process PPTX by converting to PDF first"""
        file_path = Path(file_path)
        
        if not file_path.suffix.lower() == '.pptx':
            return {
                "content": "",
                "metadata": {},
                "success": False,
                "error": "Not a PPTX file"
            }
        
        try:
            # Convert PPTX to PDF
            pdf_path = self._convert_to_pdf(file_path)
            
            if not pdf_path or not pdf_path.exists():
                return self._fallback_to_python_pptx(file_path)
            
            # Process PDF with our existing PDF processors
            from .pdf_processor import PDFProcessor
            pdf_processor = PDFProcessor(self.config)
            result = pdf_processor.process(pdf_path)
            
            # Clean up temp PDF
            try:
                pdf_path.unlink()
            except:
                pass
            
            if result.success:
                return {
                    "content": result.content,
                    "metadata": {
                        "title": result.metadata.title or file_path.stem,
                        "format": "pptx",
                        "processing_engine": "pptx-to-pdf-conversion",
                        "original_format": "pptx",
                        "converted_via": "pdf",
                        "elements": {
                            "slides": result.content.count("# Slide") if result.content else 0,
                            "tables": len(result.tables or []),
                            "images": len(result.images or []),
                            "formulas": len(result.formulas or [])
                        }
                    },
                    "images": result.images or [],
                    "tables": result.tables or [],
                    "formulas": result.formulas or [],
                    "success": True,
                    "processor": "pptx-to-pdf"
                }
            else:
                # Fallback to direct PPTX processing
                return self._fallback_to_python_pptx(file_path)
                
        except Exception as e:
            logging.warning(f"PPTX to PDF conversion failed: {e}")
            return self._fallback_to_python_pptx(file_path)
    
    def _convert_to_pdf(self, pptx_path: Path) -> Path:
        """Convert PPTX to PDF using LibreOffice"""
        temp_dir = Path(tempfile.gettempdir())
        pdf_path = temp_dir / f"{pptx_path.stem}_temp.pdf"

        auto_cleanup(pdf_path, timeout=1800)
        
        # Try LibreOffice conversion (best for PPTX)
        try:
            cmd = [
                'libreoffice', '--headless', '--convert-to', 'pdf',
                '--outdir', str(temp_dir),
                str(pptx_path)
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            # LibreOffice creates PDF with same name as input
            expected_pdf = temp_dir / f"{pptx_path.stem}.pdf"
            if expected_pdf.exists():
                expected_pdf.rename(pdf_path)
                logging.info(f"Successfully converted {pptx_path} to PDF")
                return pdf_path
            
        except Exception as e:
            logging.debug(f"LibreOffice conversion failed: {e}")
        
        # Try other methods if needed
        try:
            return self._convert_with_unoconv(pptx_path, pdf_path)
        except:
            pass
        
        logging.warning("All PDF conversion methods failed")
        return None
    
    def _convert_with_unoconv(self, pptx_path: Path, pdf_path: Path) -> Path:
        """Convert using unoconv as fallback"""
        try:
            cmd = ['unoconv', '-f', 'pdf', '-o', str(pdf_path), str(pptx_path)]
            result = subprocess.run(cmd, capture_output=True, timeout=60)
            if result.returncode == 0 and pdf_path.exists():
                return pdf_path
        except:
            pass
        return None
    
    def _fallback_to_python_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Fallback to direct python-pptx processing"""
        try:
            processor = EnhancedPptxProcessor(self.config)
            return processor.process(file_path)
        except Exception as e:
            return {
                "content": "",
                "metadata": {},
                "success": False,
                "error": f"All processing methods failed: {e}",
                "processor": "pptx-fallback"
            }


class XlsxToPdfProcessor:
    """Convert XLSX to PDF then process with PDF processors"""
    
    def __init__(self, config):
        self.config = config
    
    def process(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """Process XLSX by converting to PDF first"""
        file_path = Path(file_path)
        
        if not file_path.suffix.lower() == '.xlsx':
            return {
                "content": "",
                "metadata": {},
                "success": False,
                "error": "Not a XLSX file"
            }
        
        try:
            # Convert XLSX to PDF
            pdf_path = self._convert_to_pdf(file_path)
            
            if not pdf_path or not pdf_path.exists():
                return self._fallback_to_openpyxl(file_path)
            
            # Process PDF with our existing PDF processors
            from .pdf_processor import PDFProcessor
            pdf_processor = PDFProcessor(self.config)
            result = pdf_processor.process(pdf_path)
            
            # Clean up temp PDF
            try:
                pdf_path.unlink()
            except:
                pass
            
            if result.success:
                # Get sheet count from original XLSX for metadata
                sheet_count = self._get_sheet_count(file_path)
                
                return {
                    "content": result.content,
                    "metadata": {
                        "title": result.metadata.title or file_path.stem,
                        "format": "xlsx",
                        "processing_engine": "xlsx-to-pdf-conversion",
                        "original_format": "xlsx",
                        "converted_via": "pdf",
                        "elements": {
                            "sheets": sheet_count,
                            "tables": len(result.tables or []),
                            "images": len(result.images or []),
                            "formulas": len(result.formulas or [])
                        }
                    },
                    "images": result.images or [],
                    "tables": result.tables or [],
                    "formulas": result.formulas or [],
                    "success": True,
                    "processor": "xlsx-to-pdf"
                }
            else:
                # Fallback to direct XLSX processing
                return self._fallback_to_openpyxl(file_path)
                
        except Exception as e:
            logging.warning(f"XLSX to PDF conversion failed: {e}")
            return self._fallback_to_openpyxl(file_path)
    
    def _convert_to_pdf(self, xlsx_path: Path) -> Path:
        """Convert XLSX to PDF using available tools"""
        temp_dir = Path(tempfile.gettempdir())
        pdf_path = temp_dir / f"{xlsx_path.stem}_temp.pdf"

        auto_cleanup(pdf_path, timeout=1800)
        
        # Try multiple conversion methods in order of preference
        conversion_methods = [
            self._convert_with_libreoffice,
            self._convert_with_unoconv,
            self._convert_with_reportlab,
            self._convert_with_matplotlib
        ]
        
        for method in conversion_methods:
            try:
                if method(xlsx_path, pdf_path):
                    logging.info(f"Successfully converted {xlsx_path} to PDF using {method.__name__}")
                    return pdf_path
            except Exception as e:
                logging.debug(f"Conversion method {method.__name__} failed: {e}")
                continue
        
        logging.warning("All XLSX to PDF conversion methods failed")
        return None
    
    def _convert_with_libreoffice(self, xlsx_path: Path, pdf_path: Path) -> bool:
        """Convert using LibreOffice headless (best option)"""
        try:
            cmd = [
                'libreoffice', '--headless', '--convert-to', 'pdf',
                '--outdir', str(pdf_path.parent),
                str(xlsx_path)
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            # LibreOffice creates PDF with same name as input
            expected_pdf = pdf_path.parent / f"{xlsx_path.stem}.pdf"
            if expected_pdf.exists():
                expected_pdf.rename(pdf_path)
                return True
            
            return result.returncode == 0 and pdf_path.exists()
        except Exception as e:
            logging.debug(f"LibreOffice XLSX conversion failed: {e}")
            return False
    
    def _convert_with_unoconv(self, xlsx_path: Path, pdf_path: Path) -> bool:
        """Convert using unoconv"""
        try:
            cmd = ['unoconv', '-f', 'pdf', '-o', str(pdf_path), str(xlsx_path)]
            result = subprocess.run(cmd, capture_output=True, timeout=60)
            return result.returncode == 0 and pdf_path.exists()
        except Exception as e:
            logging.debug(f"Unoconv XLSX conversion failed: {e}")
            return False
    
    def _convert_with_reportlab(self, xlsx_path: Path, pdf_path: Path) -> bool:
        """Convert using ReportLab"""
        if not REPORTLAB_AVAILABLE or not OPENPYXL_AVAILABLE:
            return False
        
        try:
            # Read XLSX
            wb = openpyxl.load_workbook(str(xlsx_path), data_only=True)
            
            # Create PDF
            doc = SimpleDocTemplate(str(pdf_path), pagesize=A4)
            elements = []
            styles = getSampleStyleSheet()
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                
                # Add sheet title
                elements.append(Paragraph(f"Sheet: {sheet_name}", styles['Heading1']))
                elements.append(Spacer(1, 12))
                
                # Convert sheet to table data
                data = []
                for row in ws.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        data.append([str(cell) if cell is not None else "" for cell in row])
                
                if data:
                    # Limit columns and rows for PDF
                    max_cols = min(len(data[0]) if data else 0, 8)
                    max_rows = min(len(data), 50)
                    limited_data = [row[:max_cols] for row in data[:max_rows]]
                    
                    # Create table
                    if limited_data:
                        table = Table(limited_data)
                        table.setStyle(TableStyle([
                            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                            ('FONTSIZE', (0, 0), (-1, 0), 9),
                            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                            ('GRID', (0, 0), (-1, -1), 1, colors.black),
                            ('FONTSIZE', (0, 1), (-1, -1), 7)
                        ]))
                        elements.append(table)
                        elements.append(Spacer(1, 20))
                else:
                    elements.append(Paragraph("(Empty sheet)", styles['Normal']))
                    elements.append(Spacer(1, 20))
            
            wb.close()
            
            # Build PDF
            doc.build(elements)
            return pdf_path.exists()
            
        except Exception as e:
            logging.debug(f"ReportLab XLSX conversion failed: {e}")
            return False
    
    def _convert_with_matplotlib(self, xlsx_path: Path, pdf_path: Path) -> bool:
        """Convert using matplotlib (fallback)"""
        if not MATPLOTLIB_AVAILABLE or not PANDAS_AVAILABLE:
            return False
        
        try:
            # Read all sheets
            xl_file = pd.ExcelFile(str(xlsx_path))
            
            with PdfPages(str(pdf_path)) as pdf:
                for sheet_name in xl_file.sheet_names:
                    try:
                        df = pd.read_excel(xl_file, sheet_name=sheet_name)
                        
                        # Limit size for readability
                        if len(df) > 50:
                            df = df.head(50)
                        if len(df.columns) > 8:
                            df = df.iloc[:, :8]
                        
                        # Create figure
                        fig, ax = plt.subplots(figsize=(11.7, 8.27))  # A4 size
                        ax.axis('tight')
                        ax.axis('off')
                        
                        # Add title
                        fig.suptitle(f"Sheet: {sheet_name}", fontsize=14, fontweight='bold')
                        
                        # Create table
                        if not df.empty:
                            # Convert to strings and handle NaN
                            df_str = df.astype(str).replace('nan', '')
                            
                            table = ax.table(cellText=df_str.values,
                                           colLabels=df_str.columns,
                                           cellLoc='center',
                                           loc='center')
                            
                            table.auto_set_font_size(False)
                            table.set_fontsize(7)
                            table.scale(1, 1.5)
                            
                            # Style header
                            for i in range(len(df.columns)):
                                table[(0, i)].set_facecolor('#4CAF50')
                                table[(0, i)].set_text_props(weight='bold', color='white')
                        else:
                            ax.text(0.5, 0.5, 'Empty Sheet', 
                                   horizontalalignment='center', 
                                   verticalalignment='center',
                                   transform=ax.transAxes,
                                   fontsize=14)
                        
                        pdf.savefig(fig, bbox_inches='tight')
                        plt.close(fig)
                        
                    except Exception as e:
                        logging.debug(f"Failed to process sheet {sheet_name}: {e}")
                        continue
            
            return pdf_path.exists()
            
        except Exception as e:
            logging.debug(f"Matplotlib XLSX conversion failed: {e}")
            return False
    
    def _get_sheet_count(self, xlsx_path: Path) -> int:
        """Get number of sheets in XLSX file"""
        try:
            if OPENPYXL_AVAILABLE:
                wb = openpyxl.load_workbook(str(xlsx_path), read_only=True)
                count = len(wb.sheetnames)
                wb.close()
                return count
            elif PANDAS_AVAILABLE:
                xl_file = pd.ExcelFile(str(xlsx_path))
                return len(xl_file.sheet_names)
        except Exception as e:
            logging.debug(f"Failed to get sheet count: {e}")
        return 0
    
    def _fallback_to_openpyxl(self, file_path: Path) -> Dict[str, Any]:
        """Fallback to direct openpyxl processing"""
        try:
            processor = EnhancedXlsxProcessor(self.config)
            return processor.process(file_path)
        except Exception as e:
            return {
                "content": "",
                "metadata": {},
                "success": False,
                "error": f"All processing methods failed: {e}",
                "processor": "xlsx-fallback"
            }


class DoclingOfficeProcessor:
    """Advanced office processor using Docling v2.x"""
    
    def __init__(self, config: ProcessingConfig):
        self.config = config
        self.converter = None
        self._initialize()
    
    def _initialize(self):
        """Initialize Docling converter with optimized settings"""
        if not DOCLING_AVAILABLE:
            raise RuntimeError("Docling not available")
        
        try:
            format_options = {
                InputFormat.DOCX: {
                    "do_ocr": self.config.enable_ocr,
                    "do_table_structure": self.config.extract_tables,
                    "extract_images": self.config.extract_images
                },
                InputFormat.PPTX: {
                    "do_ocr": self.config.enable_ocr,
                    "do_table_structure": self.config.extract_tables,
                    "extract_images": self.config.extract_images
                },
                InputFormat.XLSX: {
                    "extract_tables": True,
                    "preserve_formulas": self.config.extract_formulas
                }
            }
            
            self.converter = DocumentConverter(format_options=format_options)
            logging.info("Docling office processor initialized successfully")
            
        except Exception as e:
            logging.error(f"Failed to initialize Docling: {e}")
            raise
    
    def process(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """Process office document with Docling"""
        try:
            result = self.converter.convert(str(file_path))
            document = result.document
            
            content_parts = []
            extracted_elements = {"text_blocks": 0, "tables": 0, "images": 0}
            
            if hasattr(document, 'export_to_markdown'):
                content = document.export_to_markdown()
                content_parts.append(content)
            
            tables = []
            if hasattr(document, 'tables') and document.tables:
                for i, table in enumerate(document.tables):
                    table_data = self._process_docling_table(table, i)
                    tables.append(table_data)
                    extracted_elements["tables"] += 1
            
            images = []
            if hasattr(document, 'pictures') and document.pictures:
                for i, picture in enumerate(document.pictures):
                    image_data = self._process_docling_image(picture, i)
                    images.append(image_data)
                    extracted_elements["images"] += 1
            
            metadata = self._extract_docling_metadata(document, file_path)
            metadata["elements"] = extracted_elements
            
            return {
                "content": "\n\n".join(content_parts),
                "metadata": metadata,
                "images": images,
                "tables": tables,
                "success": True,
                "processor": "docling"
            }
            
        except Exception as e:
            logging.error(f"Docling office processing failed for {file_path}: {e}")
            return {
                "content": "",
                "metadata": {},
                "images": [],
                "tables": [],
                "success": False,
                "error": str(e),
                "processor": "docling"
            }
    
    def _process_docling_table(self, table, index: int) -> Dict[str, Any]:
        """Process table from Docling"""
        try:
            if hasattr(table, 'export_to_markdown'):
                markdown = table.export_to_markdown()
            else:
                markdown = str(table)
            
            return {
                "index": index,
                "markdown": markdown,
                "metadata": {
                    "rows": getattr(table, 'num_rows', 0),
                    "cols": getattr(table, 'num_cols', 0)
                }
            }
        except Exception as e:
            logging.warning(f"Failed to process table {index}: {e}")
            return {"index": index, "markdown": "", "metadata": {"error": str(e)}}
    
    def _process_docling_image(self, image, index: int) -> Dict[str, Any]:
        """Process image from Docling"""
        try:
            image_info = {
                "index": index,
                "metadata": {
                    "width": getattr(image, 'width', 0),
                    "height": getattr(image, 'height', 0)
                }
            }
            
            if hasattr(image, 'get_image_data'):
                try:
                    image_data = image.get_image_data()
                    image_info["data"] = base64.b64encode(image_data).decode('utf-8')
                except:
                    pass
            
            return image_info
        except Exception as e:
            logging.warning(f"Failed to process image {index}: {e}")
            return {"index": index, "metadata": {"error": str(e)}}
    
    def _extract_docling_metadata(self, document, file_path: Path) -> Dict[str, Any]:
        """Extract metadata from Docling document"""
        metadata = {
            "title": getattr(document, 'title', file_path.stem),
            "format": file_path.suffix.lower().lstrip('.'),
            "processing_engine": "docling"
        }
        
        if hasattr(document, 'metadata'):
            doc_meta = document.metadata
            metadata.update({
                "author": getattr(doc_meta, 'author', ''),
                "created": getattr(doc_meta, 'created', ''),
                "modified": getattr(doc_meta, 'modified', '')
            })
        
        return metadata


class EnhancedDocxProcessor:
    """Enhanced DOCX processor with comprehensive extraction"""
    
    def __init__(self, config: ProcessingConfig):
        self.config = config
    
    def process(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """Process DOCX file with enhanced extraction"""
        if not PYTHON_DOCX_AVAILABLE:
            return self._create_error_result("python-docx not available")
        
        try:
            doc = DocxDocument(str(file_path))
            
            extracted_content = []
            images = []
            tables = []
            formulas = []
            comments = []
            
            metadata = self._extract_docx_metadata(doc, file_path)
            self._process_document_elements(doc, extracted_content, tables, images, formulas, comments)
            content = self._build_structured_content(extracted_content)
            
            metadata["elements"] = {
                "paragraphs": len([c for c in extracted_content if c.content_type == OfficeContentType.TEXT]),
                "tables": len(tables),
                "images": len(images),
                "formulas": len(formulas),
                "comments": len(comments)
            }
            
            return {
                "content": content,
                "metadata": metadata,
                "images": images,
                "tables": tables,
                "formulas": formulas,
                "comments": comments,
                "success": True,
                "processor": "python-docx-enhanced"
            }
            
        except Exception as e:
            logging.error(f"Enhanced DOCX processing failed: {e}")
            return self._create_error_result(str(e))
    
    def _extract_docx_metadata(self, doc: DocxDocument, file_path: Path) -> Dict[str, Any]:
        """Extract comprehensive metadata from DOCX"""
        core_props = doc.core_properties
        
        return {
            "title": core_props.title or file_path.stem,
            "author": core_props.author or "",
            "created": str(core_props.created) if core_props.created else "",
            "modified": str(core_props.modified) if core_props.modified else "",
            "subject": core_props.subject or "",
            "keywords": core_props.keywords or "",
            "format": "docx",
            "processing_engine": "python-docx-enhanced",
            "document_stats": {
                "paragraphs": len(doc.paragraphs),
                "tables": len(doc.tables),
                "sections": len(doc.sections)
            }
        }
    
    def _process_document_elements(self, doc: DocxDocument, extracted_content: List[ExtractedContent], 
                                 tables: List[Dict], images: List[Dict], 
                                 formulas: List[Dict], comments: List[Dict]):
        """Process all document elements"""
        
        for para_idx, paragraph in enumerate(doc.paragraphs):
            try:
                if paragraph.text.strip():
                    content = ExtractedContent(
                        content_type=OfficeContentType.TEXT,
                        content=paragraph.text.strip(),
                        metadata={
                            "paragraph_index": para_idx,
                            "style": paragraph.style.name if paragraph.style else "Normal"
                        },
                        formatting=self._extract_paragraph_formatting(paragraph)
                    )
                    extracted_content.append(content)
                
                hyperlinks = self._extract_hyperlinks(paragraph)
                extracted_content.extend(hyperlinks)
                
                para_images = self._extract_paragraph_images(paragraph, para_idx)
                images.extend(para_images)
                
                para_formulas = self._extract_paragraph_formulas(paragraph, para_idx)
                formulas.extend(para_formulas)
                
            except Exception as e:
                logging.warning(f"Failed to process paragraph {para_idx}: {e}")
        
        for table_idx, table in enumerate(doc.tables):
            try:
                table_data = self._process_enhanced_table(table, table_idx)
                tables.append(table_data)
                
                content = ExtractedContent(
                    content_type=OfficeContentType.TABLE,
                    content=table_data["markdown"],
                    metadata={
                        "table_index": table_idx,
                        "rows": table_data["metadata"]["rows"],
                        "cols": table_data["metadata"]["cols"]
                    }
                )
                extracted_content.append(content)
                
            except Exception as e:
                logging.warning(f"Failed to process table {table_idx}: {e}")
    
    def _extract_paragraph_formatting(self, paragraph) -> Dict[str, Any]:
        """Extract paragraph formatting"""
        formatting = {}
        try:
            if paragraph.paragraph_format:
                pf = paragraph.paragraph_format
                formatting.update({
                    "alignment": str(pf.alignment) if pf.alignment else None,
                    "left_indent": str(pf.left_indent) if pf.left_indent else None
                })
            
            if paragraph.runs:
                run = paragraph.runs[0]
                if run.font:
                    formatting["font"] = {
                        "name": run.font.name,
                        "size": str(run.font.size) if run.font.size else None,
                        "bold": run.font.bold,
                        "italic": run.font.italic
                    }
        except Exception as e:
            logging.warning(f"Failed to extract formatting: {e}")
        
        return formatting
    
    def _extract_hyperlinks(self, paragraph) -> List[ExtractedContent]:
        """Extract hyperlinks from paragraph"""
        hyperlinks = []
        try:
            for run in paragraph.runs:
                if hasattr(run._element, 'xpath'):
                    try:
                        hyperlink_elements = run._element.xpath('.//w:hyperlink')
                    except (TypeError, AttributeError):
                        # Fallback
                        hyperlink_elements = []
                        for elem in run._element.iter():
                            if 'hyperlink' in str(elem.tag):
                                hyperlink_elements.append(elem)
                    
                    for elem in hyperlink_elements:
                        try:
                            r_id = elem.get(qn('r:id'))
                            if r_id and hasattr(paragraph, 'part') and hasattr(paragraph.part, 'rels'):
                                relationship = paragraph.part.rels[r_id]
                                url = relationship.target_ref
                                text = elem.text or ""
                                
                                hyperlink = ExtractedContent(
                                    content_type=OfficeContentType.HYPERLINK,
                                    content=f"[{text}]({url})",
                                    metadata={"url": url, "text": text}
                                )
                                hyperlinks.append(hyperlink)
                        except Exception as e:
                            logging.debug(f"Failed to extract hyperlink: {e}")
        except Exception as e:
            logging.debug(f"Hyperlink extraction failed: {e}")  # Chuyển thành DEBUG
        
        return hyperlinks
    
    def _extract_paragraph_images(self, paragraph, para_idx: int) -> List[Dict[str, Any]]:
        """Extract images from paragraph"""
        images = []
        try:
            for run_idx, run in enumerate(paragraph.runs):
                try:
                    # FIX: Sử dụng xpath an toàn
                    inline_shapes = []
                    try:
                        inline_shapes = run._element.xpath('.//a:blip')
                    except (TypeError, AttributeError):
                        # Fallback method
                        for elem in run._element.iter():
                            if 'blip' in str(elem.tag):
                                inline_shapes.append(elem)
                    
                    for shape_idx, inline_shape in enumerate(inline_shapes):
                        try:
                            r_id = inline_shape.get(qn('r:embed'))
                            if r_id and hasattr(run, 'part') and hasattr(run.part, 'rels'):
                                relationship = run.part.rels[r_id]
                                image_data = relationship.target_part.blob
                                
                                image_info = {
                                    "paragraph_index": para_idx,
                                    "run_index": run_idx,
                                    "relationship_id": r_id,
                                    "format": self._detect_image_format(image_data),
                                    "size": len(image_data),
                                    "data": base64.b64encode(image_data).decode('utf-8') if self.config.extract_images else None,
                                    "metadata": {
                                        "context": paragraph.text[:100] + "..." if len(paragraph.text) > 100 else paragraph.text
                                    }
                                }
                                
                                if PIL_AVAILABLE and self.config.extract_images:
                                    try:
                                        pil_image = Image.open(io.BytesIO(image_data))
                                        image_info["metadata"].update({
                                            "width": pil_image.width,
                                            "height": pil_image.height,
                                            "mode": pil_image.mode
                                        })
                                    except Exception as e:
                                        logging.debug(f"Image analysis failed: {e}")
                                
                                images.append(image_info)
                                
                        except Exception as e:
                            logging.debug(f"Failed to extract image: {e}")
                except Exception as e:
                    logging.debug(f"Failed to process run {run_idx}: {e}")
        except Exception as e:
            logging.debug(f"Image extraction failed: {e}")  # Chuyển thành DEBUG
        
        return images
    
    def _extract_paragraph_formulas(self, paragraph, para_idx: int) -> List[Dict[str, Any]]:
        """Extract mathematical formulas"""
        formulas = []
        if not self.config.extract_formulas:
            return formulas
        
        try:
            for run in paragraph.runs:
                # FIX: Không dùng namespaces parameter
                try:
                    math_elements = run._element.xpath('.//m:oMath')
                except (TypeError, AttributeError):
                    # Fallback: tìm math elements bằng cách khác
                    math_elements = []
                    for elem in run._element.iter():
                        if 'oMath' in str(elem.tag):
                            math_elements.append(elem)
                
                for math_idx, math_elem in enumerate(math_elements):
                    try:
                        math_text = self._extract_omml_text(math_elem)
                        if math_text:
                            formulas.append({
                                "paragraph_index": para_idx,
                                "formula_index": math_idx,
                                "content": math_text,
                                "type": "omml",
                                "metadata": {"extraction_method": "omml_parsing"}
                            })
                    except Exception as e:
                        logging.debug(f"Failed to extract formula: {e}")
        except Exception as e:
            logging.debug(f"Formula extraction failed: {e}")  # Chuyển từ WARNING sang DEBUG
        
        return formulas
    
    def _extract_omml_text(self, math_element) -> str:
        """Extract text from OMML math element"""
        try:
            text_parts = []
            def extract_text_recursive(elem):
                if elem.text:
                    text_parts.append(elem.text)
                for child in elem:
                    extract_text_recursive(child)
                if elem.tail:
                    text_parts.append(elem.tail)
            
            extract_text_recursive(math_element)
            return "".join(text_parts).strip()
        except Exception as e:
            logging.debug(f"OMML text extraction failed: {e}")
            return ""
    
    def _process_enhanced_table(self, table, table_idx: int) -> Dict[str, Any]:
        """Process table with enhanced analysis"""
        try:
            rows_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                rows_data.append(row_data)
            
            if not rows_data:
                return self._create_empty_table_result(table_idx)
            
            markdown = self._table_to_enhanced_markdown(rows_data)
            table_analysis = self._analyze_table_structure(rows_data)
            
            return {
                "index": table_idx,
                "markdown": markdown,
                "raw_data": rows_data if self.config.preserve_layout else None,
                "metadata": {
                    "rows": len(rows_data),
                    "cols": len(rows_data[0]) if rows_data else 0,
                    "has_header": table_analysis["has_header"],
                    "table_type": table_analysis["table_type"],
                    "data_types": table_analysis["data_types"]
                }
            }
        except Exception as e:
            logging.warning(f"Table processing failed: {e}")
            return self._create_empty_table_result(table_idx, str(e))
    
    def _table_to_enhanced_markdown(self, rows_data: List[List[str]]) -> str:
        """Convert table to markdown"""
        if not rows_data:
            return ""
        
        try:
            max_cols = max(len(row) for row in rows_data)
            normalized_rows = []
            
            for row in rows_data:
                normalized_row = row + [""] * (max_cols - len(row))
                cleaned_row = [self._clean_cell_content(cell) for cell in normalized_row]
                normalized_rows.append(cleaned_row)
            
            markdown_rows = []
            if normalized_rows:
                header = normalized_rows[0]
                markdown_rows.append("| " + " | ".join(header) + " |")
                markdown_rows.append("| " + " | ".join(["---"] * len(header)) + " |")
                
                for row in normalized_rows[1:]:
                    markdown_rows.append("| " + " | ".join(row) + " |")
            
            return "\n".join(markdown_rows)
        except Exception as e:
            logging.warning(f"Markdown conversion failed: {e}")
            return "| Error | Converting | Table |\n| --- | --- | --- |"
    
    def _clean_cell_content(self, cell_content: str) -> str:
        """Clean cell content for markdown"""
        if not cell_content:
            return ""
        
        cleaned = cell_content.replace("\n", " ").replace("\r", " ").replace("|", "\\|")
        cleaned = " ".join(cleaned.split())
        return cleaned[:100] + "..." if len(cleaned) > 100 else cleaned
    
    def _analyze_table_structure(self, rows_data: List[List[str]]) -> Dict[str, Any]:
        """Analyze table structure"""
        if not rows_data:
            return {"has_header": False, "table_type": "empty", "data_types": []}
        
        analysis = {"has_header": False, "table_type": "data", "data_types": []}
        
        try:
            if len(rows_data) > 1:
                first_row = rows_data[0]
                second_row = rows_data[1] if len(rows_data) > 1 else []
                
                first_row_numeric = sum(1 for cell in first_row if self._is_numeric(cell))
                second_row_numeric = sum(1 for cell in second_row if self._is_numeric(cell))
                
                if first_row_numeric < len(first_row) * 0.5 and second_row_numeric > len(second_row) * 0.3:
                    analysis["has_header"] = True
            
            if len(rows_data) <= 5 and len(rows_data[0]) <= 3:
                analysis["table_type"] = "summary"
            else:
                analysis["table_type"] = "data_table"
                
        except Exception as e:
            logging.warning(f"Table analysis failed: {e}")
        
        return analysis
    
    def _is_numeric(self, value: str) -> bool:
        """Check if value is numeric"""
        if not value or not value.strip():
            return False
        value = value.strip().replace(",", "").replace("$", "").replace("%", "")
        try:
            float(value)
            return True
        except ValueError:
            return False
    
    def _detect_image_format(self, image_data: bytes) -> str:
        """Detect image format from binary data"""
        if image_data.startswith(b'\xff\xd8\xff'):
            return 'jpeg'
        elif image_data.startswith(b'\x89PNG'):
            return 'png'
        elif image_data.startswith(b'GIF8'):
            return 'gif'
        else:
            return 'unknown'
    
    def _build_structured_content(self, extracted_content: List[ExtractedContent]) -> str:
        """Build structured content from extracted elements"""
        content_parts = []
        
        for content in extracted_content:
            try:
                if content.content_type == OfficeContentType.TEXT:
                    style = content.metadata.get("style", "").lower()
                    if "heading" in style:
                        try:
                            level = int(style.split()[-1]) if style.split()[-1].isdigit() else 1
                            level = min(level, 6)
                        except:
                            level = 1
                        content_parts.append(f"\n{'#' * level} {content.content}\n")
                    else:
                        content_parts.append(content.content)
                
                elif content.content_type == OfficeContentType.TABLE:
                    content_parts.append(f"\n\n{content.content}\n\n")
                
                elif content.content_type == OfficeContentType.HYPERLINK:
                    content_parts.append(content.content)
                
            except Exception as e:
                logging.warning(f"Failed to process content element: {e}")
        
        return "\n\n".join(content_parts)
    
    def _create_error_result(self, error_msg: str) -> Dict[str, Any]:
        """Create error result"""
        return {
            "content": "",
            "metadata": {},
            "images": [],
            "tables": [],
            "formulas": [],
            "comments": [],
            "success": False,
            "error": error_msg,
            "processor": "python-docx-enhanced"
        }
    
    def _create_empty_table_result(self, table_idx: int, error: str = None) -> Dict[str, Any]:
        """Create empty table result"""
        return {
            "index": table_idx,
            "markdown": "",
            "metadata": {"rows": 0, "cols": 0, "error": error},
            "success": False
        }


class EnhancedPptxProcessor:
    """Enhanced PPTX processor with comprehensive slide analysis"""
    
    def __init__(self, config: ProcessingConfig):
        self.config = config
    
    def process(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """Process PPTX file"""
        if not PYTHON_PPTX_AVAILABLE:
            return self._create_error_result("python-pptx not available")
        
        try:
            prs = Presentation(str(file_path))
            
            content_parts = []
            images = []
            tables = []
            charts = []
            slide_metadata = []
            
            metadata = self._extract_pptx_metadata(prs, file_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_data = self._process_enhanced_slide(slide, slide_num)
                
                content_parts.append(f"\n# Slide {slide_num}\n")
                if slide_data["title"]:
                    content_parts.append(f"## {slide_data['title']}\n")
                
                content_parts.append(slide_data["content"])
                
                images.extend(slide_data["images"])
                tables.extend(slide_data["tables"])
                charts.extend(slide_data["charts"])
                slide_metadata.append(slide_data["metadata"])
            
            metadata["elements"] = {
                "slides": len(prs.slides),
                "images": len(images),
                "tables": len(tables),
                "charts": len(charts)
            }
            metadata["slide_metadata"] = slide_metadata
            
            return {
                "content": "\n\n".join(content_parts),
                "metadata": metadata,
                "images": images,
                "tables": tables,
                "charts": charts,
                "success": True,
                "processor": "python-pptx-enhanced"
            }
            
        except Exception as e:
            logging.error(f"PPTX processing failed: {e}")
            return self._create_error_result(str(e))
    
    def _extract_pptx_metadata(self, prs: Presentation, file_path: Path) -> Dict[str, Any]:
        """Extract PPTX metadata"""
        core_props = prs.core_properties
        
        metadata = {
            "title": core_props.title or file_path.stem,
            "author": core_props.author or "",
            "created": str(core_props.created) if core_props.created else "",
            "modified": str(core_props.modified) if core_props.modified else "",
            "format": "pptx",
            "processing_engine": "python-pptx-enhanced"
        }
        
        try:
            metadata["slide_dimensions"] = {
                "width": prs.slide_width,
                "height": prs.slide_height
            }
        except Exception as e:
            logging.warning(f"Failed to extract slide dimensions: {e}")
        
        return metadata
    
    def _process_enhanced_slide(self, slide, slide_num: int) -> Dict[str, Any]:
        """Process individual slide"""
        slide_data = {
            "slide_number": slide_num,
            "title": "",
            "content": "",
            "images": [],
            "tables": [],
            "charts": [],
            "metadata": {}
        }
        
        try:
            text_elements = []
            title_found = False
            
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content = shape.text.strip()
                        
                        if not title_found and (
                            hasattr(shape, 'placeholder_format') or 
                            len(text_content) < 100 and '\n' not in text_content
                        ):
                            slide_data["title"] = text_content
                            title_found = True
                        else:
                            text_elements.append(text_content)
                    
                    if self.config.extract_tables and shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table_data = self._process_pptx_table(shape.table, slide_num, len(slide_data["tables"]))
                        slide_data["tables"].append(table_data)
                    
                    elif self.config.extract_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_data = self._process_pptx_image(shape, slide_num, len(slide_data["images"]))
                        slide_data["images"].append(image_data)
                    
                    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        chart_data = self._process_pptx_chart(shape, slide_num, len(slide_data["charts"]))
                        slide_data["charts"].append(chart_data)
                
                except Exception as e:
                    logging.warning(f"Failed to process shape in slide {slide_num}: {e}")
            
            slide_data["content"] = "\n\n".join(text_elements) if text_elements else "*(No text content)*"
            
            try:
                if hasattr(slide, 'notes_slide') and slide.notes_slide:
                    notes_text = self._extract_slide_notes(slide.notes_slide)
                    if notes_text:
                        slide_data["content"] += f"\n\n**Speaker Notes:**\n{notes_text}"
            except Exception as e:
                logging.debug(f"Failed to extract notes for slide {slide_num}: {e}")
            
            slide_data["metadata"] = {
                "layout_name": self._get_slide_layout_name(slide),
                "shape_count": len(slide.shapes),
                "text_shapes": len([s for s in slide.shapes if hasattr(s, 'text') and s.text.strip()]),
                "has_title": bool(slide_data["title"]),
                "content_length": len(slide_data["content"])
            }
            
        except Exception as e:
            logging.warning(f"Failed to process slide {slide_num}: {e}")
            slide_data["metadata"]["error"] = str(e)
        
        return slide_data
    
    def _process_pptx_table(self, table, slide_num: int, table_index: int) -> Dict[str, Any]:
        """Process PPTX table"""
        try:
            rows_data = []
            for row in table.rows:
                row_data = [cell.text.strip() if cell.text else "" for cell in row.cells]
                rows_data.append(row_data)
            
            markdown = self._table_to_enhanced_markdown(rows_data)
            table_analysis = self._analyze_table_structure(rows_data)
            
            return {
                "slide": slide_num,
                "index": table_index,
                "markdown": markdown,
                "metadata": {
                    "rows": len(rows_data),
                    "cols": len(rows_data[0]) if rows_data else 0,
                    "has_header": table_analysis["has_header"],
                    "table_type": table_analysis["table_type"]
                }
            }
        except Exception as e:
            logging.warning(f"Failed to process PPTX table: {e}")
            return {
                "slide": slide_num,
                "index": table_index,
                "markdown": "",
                "metadata": {"error": str(e)},
                "success": False
            }
    
    def _process_pptx_image(self, shape, slide_num: int, image_index: int) -> Dict[str, Any]:
        """Process PPTX image"""
        try:
            image_data = {
                "slide": slide_num,
                "index": image_index,
                "shape_id": shape.shape_id,
                "name": shape.name or f"image_{image_index}",
                "metadata": {
                    "width": shape.width,
                    "height": shape.height,
                    "left": shape.left,
                    "top": shape.top
                }
            }
            
            if self.config.extract_images and hasattr(shape, 'image'):
                try:
                    image_blob = shape.image.blob
                    image_data["data"] = base64.b64encode(image_blob).decode('utf-8')
                    image_data["format"] = self._detect_image_format(image_blob)
                    image_data["size"] = len(image_blob)
                except Exception as e:
                    logging.debug(f"Failed to extract image data: {e}")
            
            return image_data
        except Exception as e:
            logging.warning(f"Failed to process PPTX image: {e}")
            return {
                "slide": slide_num,
                "index": image_index,
                "metadata": {"error": str(e)},
                "success": False
            }
    
    def _process_pptx_chart(self, shape, slide_num: int, chart_index: int) -> Dict[str, Any]:
        """Process PPTX chart"""
        try:
            chart_data = {
                "slide": slide_num,
                "index": chart_index,
                "type": "chart",
                "metadata": {"chart_type": "unknown", "has_data": False}
            }
            
            if hasattr(shape, 'chart'):
                chart = shape.chart
                chart_data["metadata"]["chart_type"] = str(chart.chart_type) if hasattr(chart, 'chart_type') else "unknown"
                
                if hasattr(chart, 'chart_title') and chart.chart_title:
                    chart_data["title"] = chart.chart_title.text_frame.text
            
            return chart_data
        except Exception as e:
            logging.warning(f"Failed to process PPTX chart: {e}")
            return {
                "slide": slide_num,
                "index": chart_index,
                "metadata": {"error": str(e)},
                "success": False
            }
    
    def _extract_slide_notes(self, notes_slide) -> str:
        """Extract speaker notes"""
        try:
            if hasattr(notes_slide, 'notes_text_frame'):
                text_frame = notes_slide.notes_text_frame
                if text_frame and hasattr(text_frame, 'text'):
                    return text_frame.text.strip()
        except Exception as e:
            logging.debug(f"Failed to extract slide notes: {e}")
        return ""
    
    def _get_slide_layout_name(self, slide) -> str:
        """Get slide layout name"""
        try:
            if hasattr(slide, 'slide_layout') and hasattr(slide.slide_layout, 'name'):
                return slide.slide_layout.name
        except:
            pass
        return "unknown"
    
    def _table_to_enhanced_markdown(self, rows_data: List[List[str]]) -> str:
        """Convert table to markdown - reuse from DOCX processor"""
        return EnhancedDocxProcessor._table_to_enhanced_markdown(self, rows_data)
    
    def _analyze_table_structure(self, rows_data: List[List[str]]) -> Dict[str, Any]:
        """Analyze table structure - reuse from DOCX processor"""
        return EnhancedDocxProcessor._analyze_table_structure(self, rows_data)
    
    def _detect_image_format(self, image_data: bytes) -> str:
        """Detect image format - reuse from DOCX processor"""
        return EnhancedDocxProcessor._detect_image_format(self, image_data)
    
    def _create_error_result(self, error_msg: str) -> Dict[str, Any]:
        """Create error result"""
        return {
            "content": "",
            "metadata": {},
            "images": [],
            "tables": [],
            "charts": [],
            "success": False,
            "error": error_msg,
            "processor": "python-pptx-enhanced"
        }


class EnhancedXlsxProcessor:
    """Enhanced XLSX processor with comprehensive data analysis"""
    
    def __init__(self, config: ProcessingConfig):
        self.config = config
    
    def process(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """Process XLSX file"""
        if not OPENPYXL_AVAILABLE:
            return self._create_error_result("openpyxl not available")
        
        try:
            workbook = openpyxl.load_workbook(str(file_path), data_only=False)
            workbook_data = openpyxl.load_workbook(str(file_path), data_only=True)
            
            content_parts = []
            tables = []
            formulas = []
            charts = []
            sheet_metadata = []
            
            metadata = self._extract_xlsx_metadata(workbook, file_path)
            
            for sheet_name in workbook.sheetnames:
                try:
                    worksheet = workbook[sheet_name]
                    worksheet_data = workbook_data[sheet_name]
                    
                    sheet_data = self._process_enhanced_worksheet(worksheet, worksheet_data, sheet_name)
                    
                    content_parts.append(f"\n# Sheet: {sheet_name}\n")
                    content_parts.append(sheet_data["content"])
                    
                    tables.extend(sheet_data["tables"])
                    formulas.extend(sheet_data["formulas"])
                    charts.extend(sheet_data["charts"])
                    sheet_metadata.append(sheet_data["metadata"])
                    
                except Exception as e:
                    logging.warning(f"Failed to process sheet {sheet_name}: {e}")
                    content_parts.append(f"\n# Sheet: {sheet_name}\n*(Error processing sheet: {e})*\n")
            
            workbook.close()
            workbook_data.close()
            
            metadata["elements"] = {
                "sheets": len(workbook.sheetnames),
                "tables": len(tables),
                "formulas": len(formulas),
                "charts": len(charts)
            }
            metadata["sheet_metadata"] = sheet_metadata
            
            return {
                "content": "\n\n".join(content_parts),
                "metadata": metadata,
                "tables": tables,
                "formulas": formulas,
                "charts": charts,
                "success": True,
                "processor": "openpyxl-enhanced"
            }
            
        except Exception as e:
            logging.error(f"XLSX processing failed: {e}")
            return self._create_error_result(str(e))
    
    def _extract_xlsx_metadata(self, workbook: Workbook, file_path: Path) -> Dict[str, Any]:
        """Extract XLSX metadata"""
        try:
            properties = workbook.properties
            
            metadata = {
                "title": properties.title or file_path.stem,
                "author": properties.creator or "",
                "created": str(properties.created) if properties.created else "",
                "modified": str(properties.modified) if properties.modified else "",
                "format": "xlsx",
                "processing_engine": "openpyxl-enhanced",
                "workbook_info": {
                    "sheet_names": workbook.sheetnames,
                    "sheet_count": len(workbook.sheetnames),
                    "active_sheet": workbook.active.title if workbook.active else ""
                }
            }
        except Exception as e:
            logging.warning(f"Failed to extract XLSX metadata: {e}")
            metadata = {
                "title": file_path.stem,
                "format": "xlsx",
                "processing_engine": "openpyxl-enhanced",
                "metadata_error": str(e)
            }
        
        return metadata
    
    def _process_enhanced_worksheet(self, worksheet: Worksheet, worksheet_data: Worksheet, 
                                  sheet_name: str) -> Dict[str, Any]:
        """Process individual worksheet"""
        sheet_data = {
            "sheet_name": sheet_name,
            "content": "",
            "tables": [],
            "formulas": [],
            "charts": [],
            "metadata": {}
        }
        
        try:
            data_range = self._find_data_range(worksheet)
            
            if data_range["max_row"] > 0 and data_range["max_col"] > 0:
                if self.config.extract_tables:
                    table_data = self._extract_worksheet_table(worksheet, worksheet_data, sheet_name, data_range)
                    sheet_data["tables"].append(table_data)
                    sheet_data["content"] = table_data["markdown"]
                
                if self.config.extract_formulas:
                    formulas = self._extract_enhanced_formulas(worksheet, sheet_name)
                    sheet_data["formulas"].extend(formulas)
                
                charts = self._extract_worksheet_charts(worksheet, sheet_name)
                sheet_data["charts"].extend(charts)
            else:
                sheet_data["content"] = "*(Empty sheet)*"
            
            sheet_data["metadata"] = {
                "dimensions": data_range,
                "data_types": self._analyze_data_types(worksheet, data_range),
                "formatting": self._analyze_formatting(worksheet, data_range)
            }
            
        except Exception as e:
            logging.warning(f"Failed to process worksheet {sheet_name}: {e}")
            sheet_data["content"] = f"*(Error processing sheet: {e})*"
            sheet_data["metadata"]["error"] = str(e)
        
        return sheet_data
    
    def _find_data_range(self, worksheet: Worksheet) -> Dict[str, int]:
        """Find actual data range in worksheet"""
        try:
            max_row = worksheet.max_row
            max_col = worksheet.max_column
            min_row = worksheet.min_row
            min_col = worksheet.min_column
            
            actual_max_row = 0
            actual_max_col = 0
            
            for row in range(min_row, max_row + 1):
                for col in range(min_col, max_col + 1):
                    cell = worksheet.cell(row=row, column=col)
                    if cell.value is not None and str(cell.value).strip():
                        actual_max_row = max(actual_max_row, row)
                        actual_max_col = max(actual_max_col, col)
            
            return {
                "min_row": min_row,
                "max_row": actual_max_row,
                "min_col": min_col,
                "max_col": actual_max_col,
                "total_rows": actual_max_row - min_row + 1 if actual_max_row > 0 else 0,
                "total_cols": actual_max_col - min_col + 1 if actual_max_col > 0 else 0
            }
        except Exception as e:
            logging.warning(f"Failed to find data range: {e}")
            return {"min_row": 1, "max_row": 0, "min_col": 1, "max_col": 0, "total_rows": 0, "total_cols": 0}
    
    def _extract_worksheet_table(self, worksheet: Worksheet, worksheet_data: Worksheet, 
                               sheet_name: str, data_range: Dict[str, int]) -> Dict[str, Any]:
        """Extract worksheet data as table"""
        try:
            if data_range["total_rows"] == 0:
                return self._create_empty_table_result(sheet_name)
            
            max_rows = min(data_range["max_row"], self.config.chunk_size or 1000)
            max_cols = min(data_range["max_col"], 100)
            
            formula_data = []
            value_data = []
            
            for row in range(data_range["min_row"], max_rows + 1):
                formula_row = []
                value_row = []
                
                for col in range(data_range["min_col"], max_cols + 1):
                    formula_cell = worksheet.cell(row=row, column=col)
                    value_cell = worksheet_data.cell(row=row, column=col)
                    
                    formula_content = str(formula_cell.value) if formula_cell.value is not None else ""
                    value_content = str(value_cell.value) if value_cell.value is not None else ""
                    
                    formula_row.append(formula_content)
                    value_row.append(value_content)
                
                formula_data.append(formula_row)
                value_data.append(value_row)
            
            display_data = value_data
            markdown = self._worksheet_to_enhanced_markdown(display_data, max_rows=50)
            table_analysis = self._analyze_spreadsheet_structure(formula_data, value_data)
            
            table_info = {
                "sheet": sheet_name,
                "markdown": markdown,
                "metadata": {
                    "dimensions": data_range,
                    "has_header": table_analysis["has_header"],
                    "data_types": table_analysis["data_types"],
                    "formula_cells": table_analysis["formula_cells"],
                    "data_quality": table_analysis["data_quality"]
                }
            }
            
            if PANDAS_AVAILABLE and len(display_data) > 1:
                try:
                    df_analysis = self._create_dataframe_analysis(display_data, table_analysis["has_header"])
                    table_info["metadata"]["dataframe_analysis"] = df_analysis
                except Exception as e:
                    logging.debug(f"DataFrame analysis failed: {e}")
            
            return table_info
        except Exception as e:
            logging.warning(f"Failed to extract worksheet table: {e}")
            return self._create_empty_table_result(sheet_name, str(e))
    
    def _worksheet_to_enhanced_markdown(self, data: List[List[str]], max_rows: int = None) -> str:
        """Convert worksheet data to markdown"""
        if not data:
            return "*(No data)*"
        
        try:
            display_data = data[:max_rows] if max_rows else data
            
            if not display_data:
                return "*(No data to display)*"
            
            if PANDAS_AVAILABLE:
                try:
                    df = pd.DataFrame(display_data[1:], columns=display_data[0] if len(display_data) > 1 else None)
                    df = df.fillna("")
                    df = df.astype(str)
                    
                    for col in df.columns:
                        df[col] = df[col].apply(lambda x: x[:50] + "..." if len(str(x)) > 50 else str(x))
                    
                    return df.to_markdown(index=False, tablefmt="pipe")
                except Exception as e:
                    logging.debug(f"Pandas markdown conversion failed: {e}")
            
            return self._manual_markdown_conversion(display_data)
        except Exception as e:
            logging.warning(f"Markdown conversion failed: {e}")
            return "*(Error converting data to markdown)*"
    
    def _manual_markdown_conversion(self, data: List[List[str]]) -> str:
        """Manual markdown conversion"""
        if not data:
            return ""
        
        markdown_rows = []
        max_cols = max(len(row) for row in data) if data else 0
        
        for i, row in enumerate(data):
            normalized_row = row + [""] * (max_cols - len(row))
            cleaned_row = [self._clean_cell_content(str(cell)) for cell in normalized_row]
            markdown_rows.append("| " + " | ".join(cleaned_row) + " |")
            
            if i == 0:
                markdown_rows.append("| " + " | ".join(["---"] * max_cols) + " |")
        
        return "\n".join(markdown_rows)
    
    def _clean_cell_content(self, content: str) -> str:
        """Clean cell content for markdown"""
        if not content:
            return ""
        
        cleaned = str(content).replace("\n", " ").replace("\r", " ").replace("|", "\\|")
        cleaned = " ".join(cleaned.split())
        return cleaned[:100] + "..." if len(cleaned) > 100 else cleaned
    
    def _analyze_spreadsheet_structure(self, formula_data: List[List[str]], 
                                     value_data: List[List[str]]) -> Dict[str, Any]:
        """Analyze spreadsheet structure"""
        analysis = {
            "has_header": False,
            "data_types": {},
            "formula_cells": 0,
            "data_quality": {}
        }
        
        try:
            if not value_data or len(value_data) < 2:
                return analysis
            
            analysis["has_header"] = self._detect_spreadsheet_header(value_data)
            
            start_row = 1 if analysis["has_header"] else 0
            num_cols = len(value_data[0]) if value_data else 0
            
            for col_idx in range(num_cols):
                col_values = [row[col_idx] for row in value_data[start_row:] if col_idx < len(row)]
                col_values = [v for v in col_values if v and str(v).strip()]
                
                if col_values:
                    col_analysis = self._analyze_column_data(col_values)
                    column_name = value_data[0][col_idx] if analysis["has_header"] and col_idx < len(value_data[0]) else f"Column_{col_idx + 1}"
                    analysis["data_types"][column_name] = col_analysis
            
            analysis["formula_cells"] = sum(
                1 for row in formula_data for cell in row 
                if cell and str(cell).startswith('=')
            )
            
            total_cells = sum(len(row) for row in value_data)
            empty_cells = sum(1 for row in value_data for cell in row if not cell or not str(cell).strip())
            
            analysis["data_quality"] = {
                "total_cells": total_cells,
                "empty_cells": empty_cells,
                "fill_rate": (total_cells - empty_cells) / total_cells if total_cells > 0 else 0,
                "has_formulas": analysis["formula_cells"] > 0
            }
            
        except Exception as e:
            logging.warning(f"Spreadsheet structure analysis failed: {e}")
            analysis["error"] = str(e)
        
        return analysis
    
    def _detect_spreadsheet_header(self, data: List[List[str]]) -> bool:
        """Detect if first row is header"""
        if len(data) < 2:
            return False
        
        try:
            first_row = data[0]
            second_row = data[1] if len(data) > 1 else []
            
            first_row_numeric = sum(1 for cell in first_row if self._is_numeric_string(str(cell)))
            second_row_numeric = sum(1 for cell in second_row if self._is_numeric_string(str(cell)))
            
            if first_row_numeric < len(first_row) * 0.3 and second_row_numeric > len(second_row) * 0.5:
                return True
        except Exception as e:
            logging.debug(f"Header detection failed: {e}")
        
        return False
    
    def _analyze_column_data(self, values: List[str]) -> Dict[str, Any]:
        """Analyze column data type"""
        analysis = {
            "primary_type": "text",
            "numeric_count": 0,
            "date_count": 0,
            "text_count": 0,
            "unique_count": 0,
            "sample_values": []
        }
        
        try:
            numeric_values = []
            
            for value in values[:100]:
                str_value = str(value).strip()
                
                if self._is_numeric_string(str_value):
                    analysis["numeric_count"] += 1
                    try:
                        numeric_values.append(float(str_value.replace(",", "")))
                    except:
                        pass
                elif self._is_date_string(str_value):
                    analysis["date_count"] += 1
                else:
                    analysis["text_count"] += 1
            
            total_analyzed = analysis["numeric_count"] + analysis["date_count"] + analysis["text_count"]
            if total_analyzed > 0:
                if analysis["numeric_count"] / total_analyzed > 0.7:
                    analysis["primary_type"] = "numeric"
                    if numeric_values:
                        analysis["numeric_stats"] = {
                            "min": min(numeric_values),
                            "max": max(numeric_values),
                            "mean": sum(numeric_values) / len(numeric_values),
                            "count": len(numeric_values)
                        }
                elif analysis["date_count"] / total_analyzed > 0.7:
                    analysis["primary_type"] = "date"
                else:
                    analysis["primary_type"] = "text"
            
            analysis["unique_count"] = len(set(values))
            analysis["sample_values"] = list(set(values))[:5]
            
        except Exception as e:
            logging.debug(f"Column analysis failed: {e}")
            analysis["error"] = str(e)
        
        return analysis
    
    def _is_numeric_string(self, value: str) -> bool:
        """Check if string is numeric"""
        if not value:
            return False
        
        cleaned = value.replace(",", "").replace("$", "").replace("%", "").strip()
        try:
            float(cleaned)
            return True
        except ValueError:
            return False
    
    def _is_date_string(self, value: str) -> bool:
        """Check if string is date"""
        if not value or len(value) < 6:
            return False
        
        date_patterns = [
            r'\d{1,2}[/-]\d{1,2}[/-]\d{2,4}',
            r'\d{4}[/-]\d{1,2}[/-]\d{1,2}',
            r'\d{1,2}-[A-Za-z]{3}-\d{2,4}',
        ]
        
        return any(re.match(pattern, value.strip()) for pattern in date_patterns)
    
    def _create_dataframe_analysis(self, data: List[List[str]], has_header: bool) -> Dict[str, Any]:
        """Create pandas DataFrame analysis"""
        try:
            if has_header and len(data) > 1:
                df = pd.DataFrame(data[1:], columns=data[0])
            else:
                df = pd.DataFrame(data)
            
            analysis = {
                "shape": df.shape,
                "columns": df.columns.tolist() if hasattr(df, 'columns') else [],
                "memory_usage": df.memory_usage(deep=True).sum() if hasattr(df, 'memory_usage') else 0
            }
            
            try:
                df_converted = df.apply(pd.to_numeric, errors='ignore')
                analysis["inferred_dtypes"] = df_converted.dtypes.astype(str).to_dict()
            except Exception as e:
                logging.debug(f"Data type inference failed: {e}")
            
            return analysis
        except Exception as e:
            logging.warning(f"DataFrame analysis failed: {e}")
            return {"error": str(e)}
    
    def _extract_enhanced_formulas(self, worksheet: Worksheet, sheet_name: str) -> List[Dict[str, Any]]:
        """Extract formulas with analysis"""
        formulas = []
        
        if not self.config.extract_formulas:
            return formulas
        
        try:
            for row in worksheet.iter_rows():
                for cell in row:
                    if cell.value and str(cell.value).startswith('='):
                        formula_info = {
                            "sheet": sheet_name,
                            "cell": cell.coordinate,
                            "formula": str(cell.value),
                            "metadata": {
                                "complexity": self._analyze_formula_complexity(str(cell.value)),
                                "references": self._extract_formula_references(str(cell.value)),
                                "functions": self._extract_formula_functions(str(cell.value))
                            }
                        }
                        
                        try:
                            if hasattr(cell, 'value'):
                                formula_info["result"] = str(cell.value)
                        except:
                            pass
                        
                        formulas.append(formula_info)
        except Exception as e:
            logging.warning(f"Formula extraction failed: {e}")
        
        return formulas
    
    def _analyze_formula_complexity(self, formula: str) -> str:
        """Analyze formula complexity"""
        try:
            if len(formula) < 20:
                return "simple"
            elif len(formula) < 100 and formula.count('(') <= 3:
                return "moderate"
            else:
                return "complex"
        except:
            return "unknown"
    
    def _extract_formula_references(self, formula: str) -> List[str]:
        """Extract cell references from formula"""
        try:
            pattern = r'[A-Z]+\d+'
            references = re.findall(pattern, formula)
            return list(set(references))
        except:
            return []
    
    def _extract_formula_functions(self, formula: str) -> List[str]:
        """Extract function names from formula"""
        try:
            pattern = r'([A-Z][A-Z0-9_]*)\s*\('
            functions = re.findall(pattern, formula.upper())
            return list(set(functions))
        except:
            return []
    
    def _extract_worksheet_charts(self, worksheet: Worksheet, sheet_name: str) -> List[Dict[str, Any]]:
        """Extract charts from worksheet"""
        charts = []
        
        try:
            if hasattr(worksheet, '_charts'):
                for i, chart in enumerate(worksheet._charts):
                    chart_info = {
                        "sheet": sheet_name,
                        "index": i,
                        "type": type(chart).__name__ if chart else "unknown",
                        "metadata": {"title": getattr(chart, 'title', '') if chart else ''}
                    }
                    charts.append(chart_info)
        except Exception as e:
            logging.debug(f"Chart extraction failed: {e}")
        
        return charts
    
    def _analyze_data_types(self, worksheet: Worksheet, data_range: Dict[str, int]) -> Dict[str, Any]:
        """Analyze data types in worksheet"""
        try:
            type_counts = {"number": 0, "text": 0, "formula": 0, "date": 0, "boolean": 0, "empty": 0}
            
            for row in range(data_range["min_row"], min(data_range["max_row"] + 1, data_range["min_row"] + 100)):
                for col in range(data_range["min_col"], min(data_range["max_col"] + 1, data_range["min_col"] + 50)):
                    cell = worksheet.cell(row=row, column=col)
                    
                    if cell.value is None:
                        type_counts["empty"] += 1
                    elif cell.data_type == 'n':
                        type_counts["number"] += 1
                    elif cell.data_type == 'f':
                        type_counts["formula"] += 1
                    elif cell.data_type == 's':
                        type_counts["text"] += 1
                    elif cell.data_type == 'd':
                        type_counts["date"] += 1
                    elif cell.data_type == 'b':
                        type_counts["boolean"] += 1
            
            return type_counts
        except Exception as e:
            logging.warning(f"Data type analysis failed: {e}")
            return {}
    
    def _analyze_formatting(self, worksheet: Worksheet, data_range: Dict[str, int]) -> Dict[str, Any]:
        """Analyze cell formatting"""
        formatting_info = {
            "has_formatting": False,
            "font_styles": set(),
            "fill_colors": set(),
            "border_styles": set()
        }
        
        try:
            sample_size = min(50, data_range["total_rows"] * data_range["total_cols"])
            cells_checked = 0
            
            for row in range(data_range["min_row"], data_range["max_row"] + 1):
                for col in range(data_range["min_col"], data_range["max_col"] + 1):
                    if cells_checked >= sample_size:
                        break
                    
                    cell = worksheet.cell(row=row, column=col)
                    
                    if cell.font and cell.font.bold:
                        formatting_info["has_formatting"] = True
                        formatting_info["font_styles"].add("bold")
                    
                    if cell.fill and cell.fill.patternType:
                        formatting_info["has_formatting"] = True
                        formatting_info["fill_colors"].add(str(cell.fill.patternType))
                    
                    cells_checked += 1
                
                if cells_checked >= sample_size:
                    break
            
            formatting_info["font_styles"] = list(formatting_info["font_styles"])
            formatting_info["fill_colors"] = list(formatting_info["fill_colors"])
            formatting_info["border_styles"] = list(formatting_info["border_styles"])
            
        except Exception as e:
            logging.debug(f"Formatting analysis failed: {e}")
        
        return formatting_info
    
    def _create_empty_table_result(self, sheet_name: str, error: str = None) -> Dict[str, Any]:
        """Create empty table result"""
        return {
            "sheet": sheet_name,
            "markdown": "*(Empty or no data)*",
            "metadata": {
                "dimensions": {"total_rows": 0, "total_cols": 0},
                "has_header": False,
                "data_types": {},
                "formula_cells": 0,
                "error": error
            },
            "success": False
        }
    
    def _create_error_result(self, error_msg: str) -> Dict[str, Any]:
        """Create error result"""
        return {
            "content": "",
            "metadata": {},
            "tables": [],
            "formulas": [],
            "charts": [],
            "success": False,
            "error": error_msg,
            "processor": "openpyxl-enhanced"
        }


# Main processor classes
class DOCXProcessor(DocumentProcessor):
    """Improved DOCX processor with PDF conversion option"""
    
    def __init__(self, config: ProcessingConfig):
        self.config = config
        self.processors = []
        super().__init__(config)
    
    def _initialize_processor(self):
        """Initialize DOCX processors with PDF conversion priority"""
        
        # Priority 1: DOCX to PDF conversion (best for formulas)
        self.processors.append(DocxToPdfProcessor(self.config))
        logging.info("DOCX to PDF converter added")
        
        # Priority 2: Enhanced python-docx (fallback)
        if PYTHON_DOCX_AVAILABLE:
            from .office_processor import EnhancedDocxProcessor
            self.processors.append(EnhancedDocxProcessor(self.config))
            logging.info("Enhanced python-docx processor added as fallback")
        
        if not self.processors:
            raise RuntimeError("No DOCX processors available")
    
    def can_process(self, file_path: Union[str, Path]) -> bool:
        return Path(file_path).suffix.lower() == '.docx'
    
    def process(self, file_path: Union[str, Path]) -> ProcessingResult:
        """Process DOCX with improved methods"""
        start_time = time.time()
        
        if not self.can_process(file_path):
            return ProcessingResult(
                content="",
                metadata=DocumentMetadata(),
                success=False,
                error_message="Not a DOCX file"
            )
        
        last_error = None
        for processor in self.processors:
            try:
                result = processor.process(file_path)
                
                if result.get("success", False):
                    processing_time = time.time() - start_time
                    
                    metadata = DocumentMetadata(
                        title=result["metadata"].get("title"),
                        author=result["metadata"].get("author"),
                        format=DocumentFormat.DOCX,
                        processing_time=processing_time,
                        extracted_elements=result["metadata"].get("elements", {})
                    )
                    
                    return ProcessingResult(
                        content=result.get("content", ""),
                        metadata=metadata,
                        images=result.get("images", []),
                        tables=result.get("tables", []),
                        formulas=result.get("formulas", []),
                        success=True
                    )
                else:
                    last_error = result.get("error", "Unknown error")
                    continue
            except Exception as e:
                last_error = str(e)
                logging.warning(f"Processor {type(processor).__name__} failed: {e}")
                continue
        
        return ProcessingResult(
            content="",
            metadata=DocumentMetadata(format=DocumentFormat.DOCX),
            success=False,
            error_message=f"All processors failed. Last error: {last_error}"
        )
    
    async def process_async(self, file_path: Union[str, Path]) -> ProcessingResult:
        loop = asyncio.get_event_loop()
        with ThreadPoolExecutor() as executor:
            return await loop.run_in_executor(executor, self.process, file_path)


class PPTXProcessor(DocumentProcessor):
    """Improved PPTX processor with PDF conversion option"""
    
    def __init__(self, config: ProcessingConfig):
        self.config = config
        self.processors = []
        super().__init__(config)
    
    def _initialize_processor(self):
        """Initialize PPTX processors with PDF conversion priority"""
        
        # Priority 1: PPTX to PDF conversion (best for layout and formulas)
        self.processors.append(PptxToPdfProcessor(self.config))
        logging.info("PPTX to PDF converter added")
        
        # Priority 2: Enhanced python-pptx (fallback)
        if PYTHON_PPTX_AVAILABLE:
            self.processors.append(EnhancedPptxProcessor(self.config))
            logging.info("Enhanced python-pptx processor added as fallback")
        
        if not self.processors:
            raise RuntimeError("No PPTX processors available")
    
    def can_process(self, file_path: Union[str, Path]) -> bool:
        return Path(file_path).suffix.lower() == '.pptx'
    
    def process(self, file_path: Union[str, Path]) -> ProcessingResult:
        """Process PPTX with improved methods"""
        start_time = time.time()
        
        if not self.can_process(file_path):
            return ProcessingResult(
                content="",
                metadata=DocumentMetadata(),
                success=False,
                error_message="Not a PPTX file"
            )
        
        last_error = None
        for processor in self.processors:
            try:
                result = processor.process(file_path)
                
                if result.get("success", False):
                    processing_time = time.time() - start_time
                    
                    metadata = DocumentMetadata(
                        title=result["metadata"].get("title"),
                        author=result["metadata"].get("author"),
                        format=DocumentFormat.PPTX,
                        processing_time=processing_time,
                        extracted_elements=result["metadata"].get("elements", {})
                    )
                    
                    return ProcessingResult(
                        content=result.get("content", ""),
                        metadata=metadata,
                        images=result.get("images", []),
                        tables=result.get("tables", []),
                        formulas=result.get("formulas", []),
                        success=True
                    )
                else:
                    last_error = result.get("error", "Unknown error")
                    continue
            except Exception as e:
                last_error = str(e)
                logging.warning(f"Processor {type(processor).__name__} failed: {e}")
                continue
        
        return ProcessingResult(
            content="",
            metadata=DocumentMetadata(format=DocumentFormat.PPTX),
            success=False,
            error_message=f"All processors failed. Last error: {last_error}"
        )
    
    async def process_async(self, file_path: Union[str, Path]) -> ProcessingResult:
        loop = asyncio.get_event_loop()
        with ThreadPoolExecutor() as executor:
            return await loop.run_in_executor(executor, self.process, file_path)


class XLSXProcessor(DocumentProcessor):
    """XLSX Document Processor with PDF conversion option"""
    
    def __init__(self, config: ProcessingConfig):
        self.config = config
        self.processors = []
        super().__init__(config)
    
    def _initialize_processor(self):
        """Initialize XLSX processors with PDF conversion priority"""
        
        # Priority 1: XLSX to PDF conversion (best for complex layouts and formulas)
        self.processors.append(XlsxToPdfProcessor(self.config))
        logging.info("XLSX to PDF converter added")
        
        # Priority 2: Docling processor (if available)
        if DOCLING_AVAILABLE:
            try:
                self.processors.append(DoclingOfficeProcessor(self.config))
                logging.info("Docling XLSX processor added as fallback")
            except Exception as e:
                logging.warning(f"Failed to initialize Docling: {e}")
        
        # Priority 3: Enhanced openpyxl processor (final fallback)
        if OPENPYXL_AVAILABLE:
            self.processors.append(EnhancedXlsxProcessor(self.config))
            logging.info("Enhanced openpyxl processor added as final fallback")
        
        if not self.processors:
            raise RuntimeError("No XLSX processors available")
    
    def can_process(self, file_path: Union[str, Path]) -> bool:
        """Check if can handle XLSX files"""
        return Path(file_path).suffix.lower() == '.xlsx'
    
    def process(self, file_path: Union[str, Path]) -> ProcessingResult:
        """Process XLSX file with improved methods"""
        start_time = time.time()
        
        if not self.can_process(file_path):
            return ProcessingResult(
                content="",
                metadata=DocumentMetadata(),
                success=False,
                error_message="Not a XLSX file"
            )
        
        last_error = None
        for processor in self.processors:
            try:
                result = processor.process(file_path)
                
                if result.get("success", False):
                    processing_time = time.time() - start_time
                    
                    metadata = DocumentMetadata(
                        title=result["metadata"].get("title"),
                        author=result["metadata"].get("author"),
                        format=DocumentFormat.XLSX,
                        processing_time=processing_time,
                        extracted_elements=result["metadata"].get("elements", {})
                    )
                    
                    return ProcessingResult(
                        content=result.get("content", ""),
                        metadata=metadata,
                        images=result.get("images", []),
                        tables=result.get("tables", []),
                        formulas=result.get("formulas", []),
                        success=True
                    )
                else:
                    last_error = result.get("error", "Unknown error")
                    continue
            except Exception as e:
                last_error = str(e)
                logging.warning(f"Processor {type(processor).__name__} failed: {e}")
                continue
        
        return ProcessingResult(
            content="",
            metadata=DocumentMetadata(format=DocumentFormat.XLSX),
            success=False,
            error_message=f"All processors failed. Last error: {last_error}"
        )
    
    async def process_async(self, file_path: Union[str, Path]) -> ProcessingResult:
        """Process XLSX asynchronously"""
        loop = asyncio.get_event_loop()
        with ThreadPoolExecutor() as executor:
            return await loop.run_in_executor(executor, self.process, file_path)


# Utility functions
def extract_office_metadata(file_path: Union[str, Path]) -> Dict[str, Any]:
    """Extract metadata without full processing"""
    file_path = Path(file_path)
    extension = file_path.suffix.lower()
    
    metadata = {
        "filename": file_path.name,
        "extension": extension,
        "size": file_path.stat().st_size,
        "modified": time.ctime(file_path.stat().st_mtime)
    }
    
    try:
        if extension == '.docx' and PYTHON_DOCX_AVAILABLE:
            doc = DocxDocument(str(file_path))
            props = doc.core_properties
            metadata.update({
                "title": props.title or "",
                "author": props.author or "",
                "created": str(props.created) if props.created else ""
            })
        elif extension == '.pptx' and PYTHON_PPTX_AVAILABLE:
            prs = Presentation(str(file_path))
            props = prs.core_properties
            metadata.update({
                "title": props.title or "",
                "author": props.author or "",
                "slide_count": len(prs.slides)
            })
        elif extension == '.xlsx' and OPENPYXL_AVAILABLE:
            wb = openpyxl.load_workbook(str(file_path), read_only=True)
            props = wb.properties
            metadata.update({
                "title": props.title or "",
                "author": props.creator or "",
                "sheet_count": len(wb.sheetnames)
            })
            wb.close()
    except Exception as e:
        metadata["extraction_error"] = str(e)
    
    return metadata


def validate_office_file(file_path: Union[str, Path]) -> Dict[str, Any]:
    """Validate office file integrity"""
    file_path = Path(file_path)
    validation = {
        "is_valid": False,
        "format": file_path.suffix.lower(),
        "errors": [],
        "warnings": []
    }
    
    try:
        if not file_path.exists():
            validation["errors"].append("File does not exist")
            return validation
        
        if file_path.stat().st_size == 0:
            validation["errors"].append("File is empty")
            return validation
        
        extension = file_path.suffix.lower()
        if extension in ['.docx', '.pptx', '.xlsx']:
            try:
                with zipfile.ZipFile(file_path, 'r') as zip_file:
                    required_files = {
                        '.docx': ['word/document.xml', '[Content_Types].xml'],
                        '.pptx': ['ppt/presentation.xml', '[Content_Types].xml'],
                        '.xlsx': ['xl/workbook.xml', '[Content_Types].xml']
                    }
                    
                    zip_contents = zip_file.namelist()
                    for required_file in required_files.get(extension, []):
                        if required_file not in zip_contents:
                            validation["errors"].append(f"Missing required file: {required_file}")
                    
                    if not validation["errors"]:
                        validation["is_valid"] = True
            except zipfile.BadZipFile:
                validation["errors"].append("File is not a valid ZIP archive")
            except Exception as e:
                validation["errors"].append(f"ZIP validation failed: {e}")
    except Exception as e:
        validation["errors"].append(f"Validation failed: {e}")
    
    return validation


# Register processors
try:
    from .base import ProcessorFactory
    ProcessorFactory.register_processor(DocumentFormat.DOCX, DOCXProcessor)
    ProcessorFactory.register_processor(DocumentFormat.PPTX, PPTXProcessor)
    ProcessorFactory.register_processor(DocumentFormat.XLSX, XLSXProcessor)
except ImportError:
    pass


# Example usage
def main():
    """Example usage"""
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python office_processor.py <office_file>")
        return
    
    file_path = Path(sys.argv[1])
    
    # Validate file
    validation = validate_office_file(file_path)
    print(f"File validation: {'✅ Valid' if validation['is_valid'] else '❌ Invalid'}")
    
    if not validation['is_valid']:
        print(f"Errors: {validation['errors']}")
        return
    
    # Extract metadata
    metadata = extract_office_metadata(file_path)
    print(f"Title: {metadata.get('title', 'N/A')}")
    print(f"Size: {metadata.get('size', 0) / 1024:.1f} KB")
    
    # Create processor
    config = ProcessingConfig(mode=ProcessingMode.BALANCED)
    
    if file_path.suffix.lower() == '.docx':
        processor = DOCXProcessor(config)
    elif file_path.suffix.lower() == '.pptx':
        processor = PPTXProcessor(config)
    elif file_path.suffix.lower() == '.xlsx':
        processor = XLSXProcessor(config)
    else:
        print("Unsupported format")
        return
    
    print(f"\nProcessing with {type(processor).__name__}...")
    
    # Process file
    result = processor.process(file_path)
    
    if result.success:
        print(f"✅ Success! Extracted {len(result.content)} characters")
        print(f"📊 Tables: {len(result.tables or [])}")
        print(f"🖼️ Images: {len(result.images or [])}")
        print(f"📐 Formulas: {len(result.formulas or [])}")
        print(f"⏱️ Processing time: {result.metadata.processing_time:.2f}s")
    else:
        print(f"❌ Failed: {result.error_message}")


if __name__ == "__main__":
    main()